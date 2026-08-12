#!/usr/bin/env python3
"""Build updated Agent Pipeline slides (Pi + skills framework)."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "figures"
OUT_DIR = Path.home() / "Desktop/DeeCamp"
NOTE = OUT_DIR / "Note"
OUT = OUT_DIR / "Agent_Pipeline_Summary_0722.pptx"


def set_run(run, text, *, size=18, bold=False, color=RGBColor(17, 17, 17)):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Arial"


def add_title(slide, text, top=0.35):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.3), Inches(0.6))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    set_run(p.add_run(), text, size=28, bold=True)
    p.alignment = PP_ALIGN.LEFT


def add_bullets(slide, lines, *, left=0.6, top=1.1, width=12, size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = 0
        p.space_after = Pt(8)
        set_run(p.add_run(), line, size=size)


def add_image_slide(prs, title, image_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_title(slide, title, top=0.2)
    # fit image under title
    pic = slide.shapes.add_picture(
        str(image_path),
        Inches(0.35),
        Inches(0.85),
        width=Inches(12.6),
    )
    # if too tall, scale by height
    max_h = Inches(6.3)
    if pic.height > max_h:
        ratio = max_h / pic.height
        pic.height = max_h
        pic.width = int(pic.width * ratio)
        pic.left = int((prs.slide_width - pic.width) / 2)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 Title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Pre-FM Agent Pipeline", top=2.4)
    box = s.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(12), Inches(1.5))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    set_run(p.add_run(),
            "Pi coding agent · skills + scripts · scATAC / Multiome → FM-ready corpus",
            size=18)
    p2 = tf.add_paragraph()
    set_run(p2.add_run(), "Updated 2026-07-22", size=14, color=RGBColor(80, 80, 80))

    # 2 Scope
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Scope and Positioning")
    add_bullets(s, [
        "Goal: turn public scATAC-seq / multiome data into a reproducible, auditable, FM-ready corpus",
        "Difference from EpiAgent: agent-assisted curation + multiome + standardized handoff (not manual-only scATAC)",
        "Principle: Agent orchestrates trusted tools; it does not invent bioinformatics algorithms",
        "Harness: Pi coding agent (same pattern as scIsoAgent) — skills/ + scripts/ as siblings",
    ])

    # 3 Agent Framework figure
    add_image_slide(prs, "Agent Framework", FIG / "agent_framework.png")

    # 4 Framework layers text
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Agent Framework — Layers")
    add_bullets(s, [
        "Request Interpretation: curation / processing / handoff / continue from results",
        "Skill Selection: entry pipelines + analysis leaves + corpus leaves (13 skills)",
        "Execution Control: validate → resolve paths → show command → confirm → --stage launch",
        "Session Management: Pi shell · conda envs · provenance · resume",
        "Execution Environment: skills/ + scripts/ · SnapATAC2 / scanpy / muon",
        "Outputs: catalog/QC · h5ad/h5mu · cell×cCRE · data cards · FM corpus",
    ], size=17)

    # 5 Skills inventory
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Skills Inventory (13)")
    add_bullets(s, [
        "Router (1): sc-epi-agent",
        "Entry (3): curation-pipeline · processing-pipeline · handoff-pipeline",
        "Analysis (5): scatac-fragment-qc · scatac-peak-matrix · scrna-qc · multiome-qc · map-to-ccre",
        "Corpus (4): resource-setup · download-validate · tokenize-cell-sentence · fm-handoff",
        "Contract: SKILL.md = when / params / QC / failure / human-review; scripts/*.py = execution",
        "Run: ./setup_pi.sh → pi → /skill:<name>",
    ], size=17)

    # 6 Workflow figure
    add_image_slide(prs, "Workflow Overview", FIG / "workflow_overview.png")

    # 7 Workflow stages text
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Workflow — Three Stages")
    add_bullets(s, [
        "1. Inputs & Curation: discover → metadata → manifest → eligibility → cCRE vocab → download",
        "2. Modality Processing: 7A SnapATAC2 (preferred) · 7B peak-matrix · 7C Multiome · 7D scRNA-ref",
        "3. Representation & Handoff: map-to-ccre → TF-IDF / cell sentence → data cards → FM corpus",
        "QC is stage-aware: plot → confirm thresholds → filter → embed → cluster → doublet → finalize",
    ], size=17)

    # 8 Processing branches
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Modality Branches")
    add_bullets(s, [
        "7A fragments (preferred): SnapATAC2 import → TSSe → QC → spectral → Leiden → doublet → cCRE",
        "7B peak matrix (fallback): approximate peak→cCRE mapping; flag representation_quality",
        "7C Multiome: barcode pair-check → RNA (scanpy) + ATAC (SnapATAC2) QC → paired-pass ∩",
        "7D scRNA reference: scanpy only; supports label transfer — not in ATAC pretraining corpus",
    ], size=17)

    # 9 Unified representation
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Unified Representation & Handoff")
    add_bullets(s, [
        "All ATAC datasets → shared cCRE feature space on GRCh38",
        "Outputs: sparse cell×cCRE · TF-IDF / LSI · TF-IDF-ranked cell sentences",
        "Optional cell×gene RNA (multiome / reference)",
        "Handoff: vocabulary + matrices + tokens + train/val/test + data cards + provenance",
    ])

    # 10 MVP
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "MVP and Open Decisions")
    add_bullets(s, [
        "MVP: 10x PBMC scATAC · 10x PBMC Multiome · hematopoiesis / brain scATAC",
        "P0: Pi mount + fill scatac_fragment_qc / scrna_qc stages",
        "P1: multiome_qc + map_to_ccre + cCRE vocab",
        "P2: tokenize + fm-handoff + controlled download",
        "Open: feature space choice · QC fixed vs adaptive · global vs per-dataset TF-IDF",
    ], size=17)

    # 11 References
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Key References")
    add_bullets(s, [
        "Chen et al. EpiAgent. Nature Methods, 2025.",
        "Stuart et al. Signac. Nature Methods, 2021.",
        "Granja et al. ArchR. Nature Genetics, 2021.",
        "ENCODE Project Consortium. Nature, 2020.",
        "Heumos et al. Best practices for single-cell analysis. Nat Rev Genetics, 2023.",
        "scIsoAgent / Pi coding agent — skills + scripts portable bundle pattern.",
    ], size=16)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    # also copy next to Note
    prs.save(NOTE / "Agent_Pipeline_Summary_0722.pptx")
    print("wrote", OUT)


if __name__ == "__main__":
    main()
